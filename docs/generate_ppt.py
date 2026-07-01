#!/usr/bin/env python3
"""生成智能刷题系统课程设计答辩PPT。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── 配色 ───
BLUE = RGBColor(0x0E, 0xA5, 0xE9)
GREEN = RGBColor(0x10, 0xB9, 0x81)
DARK = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF0, 0xF9, 0xFF)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)


def add_bg(slide, color=LIGHT_BG):
    """设置幻灯片背景色。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_box(slide, text, left=Inches(0.5), top=Inches(0.3),
                  width=Inches(9), height=Inches(0.8), size=Pt(32),
                  color=DARK, bold=True):
    """添加标题文本框。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = '微软雅黑'
    return txBox


def add_text_box(slide, text, left, top, width, height,
                 size=Pt(16), color=DARK, bold=False, align=PP_ALIGN.LEFT):
    """添加文本框。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = '微软雅黑'
    p.alignment = align
    return txBox


def add_bullet_list(slide, items, left, top, width, height,
                    size=Pt(16), color=DARK, bullet_color=BLUE):
    """添加带圆点的列表。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = size
        p.font.color.rgb = color
        p.font.name = '微软雅黑'
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_table(slide, headers, rows, left, top, width, height):
    """添加表格。"""
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = '微软雅黑'
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    # 数据
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK
                p.font.name = '微软雅黑'
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)

    return table_shape


def add_code_box(slide, code, left, top, width, height):
    """添加代码框（深色背景）。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0xA6, 0xE3, 0xA1)
    p.font.name = 'Consolas'
    return shape


def add_accent_bar(slide, left=Inches(0.5), top=Inches(1.1), width=Inches(1.5)):
    """添加装饰条。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()


def add_stat_card(slide, label, value, left, top, color=BLUE):
    """添加数据卡片。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = '微软雅黑'

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    p2.font.name = '微软雅黑'
    p2.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # ═══════════════════════════════════════════════
    #  第1页：封面
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白
    add_bg(slide, DARK)

    # 装饰条
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Pt(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

    add_text_box(slide, '智能刷题系统', Inches(1), Inches(1.2), Inches(8), Inches(0.8),
                 size=Pt(40), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 'Smart Quiz System', Inches(1), Inches(1.9), Inches(8), Inches(0.5),
                 size=Pt(20), color=BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 'Python 程序设计课程设计', Inches(1), Inches(2.5), Inches(8), Inches(0.4),
                 size=Pt(16), color=GRAY, align=PP_ALIGN.CENTER)

    info = '姓名：jxh    班级：24软件工程1班    指导教师：***\n信息工程学院    2026年6月'
    add_text_box(slide, info, Inches(1), Inches(3.8), Inches(8), Inches(0.8),
                 size=Pt(14), color=GRAY, align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════
    #  第2页：我发现了什么问题？
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_box(slide, '01  我发现了什么问题？')
    add_accent_bar(slide)

    add_text_box(slide, '真实场景：大学生备考《Python程序设计》等课程时，需要大量刷题巩固知识点',
                 Inches(0.5), Inches(1.3), Inches(9), Inches(0.4), size=Pt(14), color=GRAY)

    add_table(slide,
        ['问题', '具体表现'],
        [
            ['📄 题目管理混乱', '从学习通复制的题目格式杂乱，包含"教师批阅""得X分"等干扰信息'],
            ['❌ 错题无法汇总', '做错的题目散落在各次练习中，无法集中复习薄弱环节'],
            ['📊 学习效果难以量化', '不知道哪些题型弱、正确率如何变化，复习缺乏针对性'],
            ['🔄 复习策略不科学', '没有依据遗忘曲线安排复习，考前突击效果差、遗忘快'],
        ],
        Inches(0.5), Inches(1.8), Inches(9), Inches(2.8))

    # ═══════════════════════════════════════════════
    #  第3页：为什么值得解决？
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_box(slide, '02  为什么值得解决？')
    add_accent_bar(slide)

    add_bullet_list(slide, [
        '大学生平均每周花费 3-5 小时进行课程复习刷题',
        '其中约 30% 的时间浪费在题目整理和格式清洗上',
        '70% 的学生没有系统的错题管理习惯',
        '艾宾浩斯遗忘曲线：学习后第1/2/4/7/15/30天是最佳复习时间点',
    ], Inches(0.5), Inches(1.3), Inches(4.5), Inches(2.5))

    add_table(slide,
        ['影响', '说明'],
        [
            ['⏱️ 时间浪费', '手动整理题目格式，重复劳动'],
            ['📉 复习低效', '没有针对性，不知道薄弱点在哪'],
            ['🧠 记忆衰减', '缺乏间隔复习策略，考前突击遗忘快'],
        ],
        Inches(5.2), Inches(1.3), Inches(4.3), Inches(1.8))

    add_text_box(slide, '结论：一个能自动解析题目、智能管理错题、科学安排复习、量化学习效果的工具，具有明确的实际需求。',
                 Inches(0.5), Inches(4.2), Inches(9), Inches(0.5), size=Pt(14), bold=True, color=BLUE)

    # ═══════════════════════════════════════════════
    #  第4页：项目目标
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_box(slide, '03  项目目标')
    add_accent_bar(slide)

    add_table(slide,
        ['核心问题', '对应目标', '技术方案'],
        [
            ['题目整理耗时', '一键导入，自动解析', '正则+AI双层解析引擎'],
            ['错题无法汇总', '自动记录，智能推荐', '答题记录+错题排序'],
            ['复习策略不科学', '间隔复习', '艾宾浩斯曲线推送'],
            ['学习效果不明确', '可视化报表', 'Plotly交互式图表'],
        ],
        Inches(0.5), Inches(1.3), Inches(9), Inches(2.2))

    add_stat_card(slide, '节省时间', '100题/3秒', Inches(0.5), Inches(3.8), BLUE)
    add_stat_card(slide, '提升效率', '错题+间隔', Inches(2.7), Inches(3.8), GREEN)
    add_stat_card(slide, '数据驱动', '可视化', Inches(4.9), Inches(3.8), PURPLE)
    add_stat_card(slide, '安全可靠', 'PBKDF2', Inches(7.1), Inches(3.8), DARK)

    # ═══════════════════════════════════════════════
    #  第5页：智能题目导入
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_box(slide, '04  核心功能 — 智能题目导入')
    add_accent_bar(slide)

    add_bullet_list(slide, [
        '四种来源：文本粘贴 / PDF / DOCX / 图片OCR',
        '双层解析引擎：正则高速（1000题/秒）+ AI智能兜底（DeepSeek/MIMO）',
        '置信度路由：≥0.8直接采纳 → 0.5~0.8 AI修复 → <0.5 AI重解析',
        '智能清洗：全角转半角、过滤"教师批阅""得X分""AI讲解"等噪音',
        '语义指纹去重（MD5）+ 质量检查 + 自动修复',
    ], Inches(0.5), Inches(1.3), Inches(5), Inches(3))

    add_code_box(slide,
        '输入：1.(单选题) 关于视图描述不正确的是（ ）\n'
        '      A.视图从表导出  B.视图存储实际数据\n'
        '      我的答案：B; 得2分; 教师批阅时间\n\n'
        '输出：题型=单选 | 答案=B ✅\n'
        '      置信度=0.92 → 直接采纳',
        Inches(5.7), Inches(1.3), Inches(3.8), Inches(3.2))

    # ═══════════════════════════════════════════════
    #  第6页：多模式刷题
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_box(slide, '04  核心功能 — 多模式刷题')
    add_accent_bar(slide)

    add_table(slide,
        ['模式', '说明', '适用场景'],
        [
            ['🎲 随机模式', '从题库随机抽题', '日常练习、综合复习'],
            ['❌ 错题重练', '按错误次数排序', '薄弱知识点强化'],
            ['⭐ 难度专项', '按难度筛选', '针对性突破'],
            ['📅 间隔复习', '艾宾浩斯曲线', '科学长期记忆'],
        ],
        Inches(0.5), Inches(1.3), Inches(5.5), Inches(2.2))

    add_bullet_list(slide, [
        '支持按题型和学科筛选',
        '实时计时，即时反馈对错',
        '练习进度持久化（中断可恢复）',
        '上一题/下一题导航',
        '答案判断5层匹配策略',
    ], Inches(6.2), Inches(1.3), Inches(3.3), Inches(3))

    # ═══════════════════════════════════════════════
    #  第7页：模拟考试
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_box(slide, '04  核心功能 — 模拟考试')
    add_accent_bar(slide)

    add_bullet_list(slide, [
        '自定义组卷：题数、时长、题型、学科',
        '限时考试：实时倒计时 + 进度条',
        '低时间警告：剩余 <5 分钟视觉提醒',
        '超时自动交卷：倒计时归零后自动触发交卷',
        '成绩报告：总分、正确率、用时、逐题详情',
    ], Inches(0.5), Inches(1.3), Inches(4.5), Inches(3))

    add_code_box(slide,
        '设置参数 → 随机抽题\n'
        '    → 进入考试（倒计时）\n'
        '    → 逐题作答\n'
        '    → <5分钟警告\n'
        '    → ≤0 自动交卷\n'
        '    → 逐题评分\n'
        '    → 成绩报告',
        Inches(5.5), Inches(1.3), Inches(4), Inches(3.2))

    # ═══════════════════════════════════════════════
    #  第8页：学习统计
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_box(slide, '04  核心功能 — 学习统计与可视化')
    add_accent_bar(slide)

    add_bullet_list(slide, [
        '仪表盘：题库总量、答题次数、正确率一目了然',
        '趋势图：每日答题柱状图 + 正确率折线图（Plotly交互式）',
        '题型分析：各题型正确率对比，精准定位薄弱环节',
        '智能建议：基于数据自动生成个性化学习建议',
        '时间范围筛选：近7天/14天/30天/90天/全部',
        '答题历史：支持CSV导出（UTF-8 BOM兼容Excel）',
    ], Inches(0.5), Inches(1.3), Inches(5.5), Inches(3.5))

    add_text_box(slide, '📊 示例：\n系统检测到你的"多选题"\n正确率仅 45%，\n建议专项练习。',
                 Inches(6.5), Inches(1.5), Inches(3), Inches(2), size=Pt(14), color=BLUE)

    # ═══════════════════════════════════════════════
    #  第9页：为什么选择Python
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_box(slide, '05  为什么选择 Python？')
    add_accent_bar(slide)

    add_table(slide,
        ['能力', 'Python方案', '为什么选它'],
        [
            ['📄 文件处理', 'pdfplumber + python-docx', '生态成熟，几行代码解析'],
            ['📊 数据分析', 'pandas + SQLite3', '内置sqlite3，无需安装'],
            ['📈 数据可视化', 'Plotly', '交互式图表，嵌入Web'],
            ['🌐 Web开发', 'Streamlit', '纯Python，无需前端'],
            ['🤖 AI扩展', 'OpenAI SDK', '兼容DeepSeek/MIMO'],
            ['🔐 密码安全', 'hashlib PBKDF2', '100,000次迭代'],
            ['🧪 测试', 'pytest', '75个用例，全部通过'],
        ],
        Inches(0.5), Inches(1.3), Inches(9), Inches(3.2))

    add_text_box(slide, '一门语言覆盖数据处理、文件解析、Web界面、可视化、AI全链路',
                 Inches(0.5), Inches(4.7), Inches(9), Inches(0.4), size=Pt(14), bold=True, color=BLUE)

    # ═══════════════════════════════════════════════
    #  第10页：系统架构
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_box(slide, '06  系统架构')
    add_accent_bar(slide)

    layers = [
        ('用户交互层', '9个页面 + 主题CSS + 公共组件', BLUE),
        ('业务逻辑层', '导入引擎 + 练习管理 + AI增强', GREEN),
        ('数据处理层', '正则解析（20+预编译）+ AI解析', PURPLE),
        ('数据存储层', 'SQLite单例管理器（WAL+线程安全）', DARK),
    ]
    for i, (name, desc, color) in enumerate(layers):
        y = Inches(1.3) + Inches(i * 0.85)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.5), y, Inches(5.5), Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f'{name}  —  {desc}'
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = '微软雅黑'

    add_bullet_list(slide, [
        '单例模式：数据库全局唯一连接，双重检查锁',
        '策略模式：正则解析与AI解析可自由切换',
        '模块化插拔：新增题型/页面只需添加模块',
    ], Inches(6.3), Inches(1.3), Inches(3.2), Inches(2.5))

    # ═══════════════════════════════════════════════
    #  第11页：关键技术
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_box(slide, '07  关键技术实现')
    add_accent_bar(slide)

    add_text_box(slide, 'PBKDF2 密码安全（100,000次迭代 + 32字符随机盐）',
                 Inches(0.5), Inches(1.2), Inches(4.5), Inches(0.3), size=Pt(13), bold=True)
    add_code_box(slide,
        'def _hash_password(password, salt=None):\n'
        '    iterations = 100_000\n'
        '    salt = salt or secrets.token_hex(16)\n'
        '    dk = hashlib.pbkdf2_hmac(\n'
        '        "sha256", password.encode(),\n'
        '        salt.encode(), iterations)\n'
        '    return f"pbkdf2:{iterations}:{salt}:{dk.hex()}"',
        Inches(0.5), Inches(1.6), Inches(4.5), Inches(2.2))

    add_text_box(slide, '统一答案判断（5层匹配策略）',
                 Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3), size=Pt(13), bold=True)
    add_code_box(slide,
        'def check_answer(user, correct, q_type):\n'
        '    if q_type in ("fill", "short"):\n'
        '        if user == correct: return True    # 1.精确\n'
        '        if clean(u)==clean(c): return True  # 2.去括号\n'
        '        if correct in user: return True     # 3.包含\n'
        '        if parts_match(u,c): return True    # 4.子集\n'
        '        if fuzzy_match(u,c): return True    # 5.模糊\n'
        '    return normalize(u) == normalize(c)',
        Inches(5.2), Inches(1.6), Inches(4.3), Inches(2.2))

    add_text_box(slide, '核心知识点：单例模式 · 策略模式 · 正则预编译 · PBKDF2 · SQLite WAL · Streamlit缓存',
                 Inches(0.5), Inches(4.2), Inches(9), Inches(0.4), size=Pt(12), color=GRAY)

    # ═══════════════════════════════════════════════
    #  第12页：项目成果
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_box(slide, '08  项目成果')
    add_accent_bar(slide)

    add_table(slide,
        ['功能', '状态', '说明'],
        [
            ['智能题目导入', '✅', '文本/PDF/DOCX/OCR，双层解析'],
            ['多模式刷题', '✅', '随机/错题/难度/间隔复习'],
            ['模拟考试', '✅', '限时考试，自动交卷'],
            ['学习统计', '✅', 'Plotly可视化+智能建议'],
            ['用户系统', '✅', 'PBKDF2加密+数据隔离'],
            ['间隔复习', '✅', '艾宾浩斯曲线（1/2/4/7/15/30天）'],
            ['测试覆盖', '✅', '75个单元测试，全部通过'],
        ],
        Inches(0.5), Inches(1.3), Inches(5.5), Inches(3.3))

    add_table(slide,
        ['指标', '改进前', '改进后'],
        [
            ['题目导入', '手动整理数小时', '3秒自动解析'],
            ['错题复习', '手动翻找', '自动汇总排序'],
            ['学习追踪', '凭感觉', '可视化报表'],
            ['复习策略', '考前突击', '间隔复习'],
            ['数据安全', '无', 'PBKDF2+隔离'],
        ],
        Inches(6.2), Inches(1.3), Inches(3.3), Inches(2.8))

    # ═══════════════════════════════════════════════
    #  第13页：测试
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_box(slide, '09  测试与质量保障')
    add_accent_bar(slide)

    add_stat_card(slide, '测试用例', '75', Inches(0.5), Inches(1.3), BLUE)
    add_stat_card(slide, '通过率', '100%', Inches(2.7), Inches(1.3), GREEN)
    add_stat_card(slide, '执行时间', '1.89s', Inches(4.9), Inches(1.3), PURPLE)
    add_stat_card(slide, '已修复问题', '20个', Inches(7.1), Inches(1.3), DARK)

    add_table(slide,
        ['测试类别', '用例数', '覆盖内容'],
        [
            ['数据库', '4', '表创建/外键/收藏/导入历史'],
            ['解析器', '4', '单选/多选/填空/简答'],
            ['答案判断', '8', '精确/括号/包含/模糊/多选/判断'],
            ['AI辅助', '7', 'JSON提取/规范化/标准化'],
            ['密码安全', '6', 'PBKDF2/兼容/随机盐'],
            ['边界条件', '8', '空文本/超长/空答案'],
        ],
        Inches(0.5), Inches(2.8), Inches(9), Inches(2.3))

    # ═══════════════════════════════════════════════
    #  第14页：未来展望
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_box(slide, '10  未来展望')
    add_accent_bar(slide)

    add_text_box(slide, '已实现的"未来功能"：', Inches(0.5), Inches(1.3), Inches(4), Inches(0.3),
                 size=Pt(16), bold=True)
    add_bullet_list(slide, [
        '✅ 用户登录系统 → 已实现注册/登录 + PBKDF2',
        '✅ 多用户数据隔离 → 已实现 user_id 字段绑定',
        '✅ 移动端适配 → 已实现响应式布局',
    ], Inches(0.5), Inches(1.7), Inches(4.5), Inches(1.5))

    add_text_box(slide, '后续迭代方向：', Inches(5.5), Inches(1.3), Inches(4), Inches(0.3),
                 size=Pt(16), bold=True)
    add_bullet_list(slide, [
        '🤖 本地AI模型：Tesseract OCR离线识别',
        '🏫 多学科题库：覆盖更多课程',
        '☁️ 云端同步：数据多设备同步',
        '📊 学习报告：PDF格式自动生成',
    ], Inches(5.5), Inches(1.7), Inches(4), Inches(2))

    # ═══════════════════════════════════════════════
    #  第15页：感谢
    # ═══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Pt(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

    add_text_box(slide, '谢谢大家！', Inches(1), Inches(1.2), Inches(8), Inches(0.8),
                 size=Pt(40), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, '智能刷题系统', Inches(1), Inches(2), Inches(8), Inches(0.5),
                 size=Pt(24), color=BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, '让学习更高效，让复习更有针对性', Inches(1), Inches(2.6), Inches(8), Inches(0.4),
                 size=Pt(16), color=GRAY, align=PP_ALIGN.CENTER)

    add_text_box(slide, '75个测试 · 5种题型 · 4种练习模式 · 6张数据表 · 20个已修复问题',
                 Inches(1), Inches(3.5), Inches(8), Inches(0.4),
                 size=Pt(14), color=GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, '姓名：jxh    班级：24软件工程1班',
                 Inches(1), Inches(4.2), Inches(8), Inches(0.4),
                 size=Pt(14), color=GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, '欢迎提问！ Q & A', Inches(1), Inches(4.8), Inches(8), Inches(0.4),
                 size=Pt(18), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # ─── 保存 ───
    output = '智能刷题系统_答辩PPT_jxh_v2.pptx'
    prs.save(output)
    print(f'✅ PPT已生成: {output}')


if __name__ == '__main__':
    main()
